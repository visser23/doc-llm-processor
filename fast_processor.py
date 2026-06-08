#!/usr/bin/env python3
"""
Fast Document Processing Pipeline - Native Extraction First
Processes PDFs 100-1000x faster than OCR-based approaches by using native text extraction.
Falls back to OCR only for truly scanned documents.
"""

import os
import json
import argparse
import hashlib
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from concurrent.futures import ProcessPoolExecutor, as_completed
from dataclasses import dataclass, asdict
import multiprocessing

# Third-party imports
import fitz  # PyMuPDF - fast native PDF extraction
import pandas as pd
from tqdm import tqdm

# Optional: python-docx for Word documents
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# Optional: python-pptx for PowerPoint
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@dataclass
class ProcessingStats:
    """Statistics for processing run."""
    total_files: int = 0
    successful: int = 0
    failed: int = 0
    total_chunks: int = 0
    total_pages: int = 0
    processing_time_seconds: float = 0.0
    files_per_second: float = 0.0
    used_ocr_fallback: int = 0


class FastDocumentProcessor:
    """
    High-performance document processor using native text extraction.
    
    Key features:
    - PyMuPDF for direct PDF text extraction (no OCR for born-digital PDFs)
    - Sentence-aware chunking to avoid breaking mid-thought
    - Parallel processing for multiple files
    - Automatic OCR fallback for scanned documents
    """
    
    def __init__(self,
                 input_dir: str,
                 output_dir: str = "processed_knowledge_base",
                 chunk_size: int = 1000,
                 overlap: int = 200,
                 min_text_density: float = 0.01,
                 workers: int = None,
                 enable_ocr_fallback: bool = True):
        """
        Initialize the fast document processor.
        
        Args:
            input_dir: Directory containing documents to process
            output_dir: Directory to save processed output
            chunk_size: Target characters per chunk
            overlap: Character overlap between chunks
            min_text_density: Minimum text/page ratio before triggering OCR fallback
            workers: Number of parallel workers (default: CPU count)
            enable_ocr_fallback: Whether to use OCR for scanned documents
        """
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.min_text_density = min_text_density
        self.workers = workers or max(1, multiprocessing.cpu_count() - 1)
        self.enable_ocr_fallback = enable_ocr_fallback
        
        # Create output directory
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Supported file extensions
        self.supported_extensions = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.md'}
        
        # Track results
        self.stats = ProcessingStats()
        self.failed_files: List[Dict[str, str]] = []
        
    def get_file_hash(self, file_path: Path) -> str:
        """Generate MD5 hash for file tracking."""
        with open(file_path, 'rb') as f:
            return hashlib.md5(f.read()).hexdigest()
    
    def find_documents(self) -> List[Path]:
        """Find all supported documents in input directory."""
        documents = []
        for ext in self.supported_extensions:
            documents.extend(self.input_dir.rglob(f"*{ext}"))
        return sorted(documents)
    
    # =========================================================================
    # NATIVE TEXT EXTRACTION (FAST PATH)
    # =========================================================================
    
    def extract_pdf_native(self, file_path: Path) -> Tuple[List[Dict[str, Any]], bool]:
        """
        Extract text directly from PDF using PyMuPDF.
        
        Returns:
            Tuple of (pages_data, needs_ocr_fallback)
        """
        pages = []
        total_text_length = 0
        total_area = 0
        
        try:
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Get page dimensions for density calculation
                rect = page.rect
                page_area = rect.width * rect.height
                total_area += page_area
                
                # Extract text - this is the FAST operation
                text = page.get_text("text")
                total_text_length += len(text)
                
                if text.strip():
                    pages.append({
                        "page_number": page_num + 1,
                        "text": text,
                        "char_count": len(text)
                    })
            
            doc.close()
            
            # Check if document appears to be scanned (very low text density)
            text_density = total_text_length / max(total_area, 1) if total_area > 0 else 0
            needs_ocr = len(pages) == 0 or (text_density < self.min_text_density and total_text_length < 100)
            
            return pages, needs_ocr
            
        except Exception as e:
            print(f"  Error in native extraction: {e}")
            return [], True
    
    def extract_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from Word documents."""
        if not HAS_DOCX:
            print(f"  Warning: python-docx not installed, skipping {file_path.name}")
            return []
        
        try:
            doc = DocxDocument(file_path)
            pages = []
            current_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    current_text.append(para.text)
            
            # Word docs don't have real pages, treat as single page
            full_text = "\n".join(current_text)
            if full_text.strip():
                pages.append({
                    "page_number": 1,
                    "text": full_text,
                    "char_count": len(full_text)
                })
            
            return pages
        except Exception as e:
            print(f"  Error extracting DOCX: {e}")
            return []
    
    def extract_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint presentations."""
        if not HAS_PPTX:
            print(f"  Warning: python-pptx not installed, skipping {file_path.name}")
            return []
        
        try:
            prs = Presentation(file_path)
            pages = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text = "\n".join(slide_text)
                    pages.append({
                        "page_number": slide_num,
                        "text": text,
                        "char_count": len(text)
                    })
            
            return pages
        except Exception as e:
            print(f"  Error extracting PPTX: {e}")
            return []
    
    def extract_text_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            
            if text.strip():
                return [{
                    "page_number": 1,
                    "text": text,
                    "char_count": len(text)
                }]
            return []
        except Exception as e:
            print(f"  Error reading text file: {e}")
            return []
    
    # =========================================================================
    # SMART CHUNKING
    # =========================================================================
    
    def smart_chunk_text(self, text: str, source_page: int = None) -> List[Dict[str, Any]]:
        """
        Chunk text intelligently, respecting sentence boundaries.
        
        Args:
            text: Text to chunk
            source_page: Original page number for metadata
            
        Returns:
            List of chunk dictionaries
        """
        if not text or not text.strip():
            return []
        
        # Sentence splitting pattern (handles common abbreviations)
        sentence_pattern = r'(?<=[.!?])\s+(?=[A-Z])'
        sentences = re.split(sentence_pattern, text)
        
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_length = len(sentence)
            
            # If single sentence exceeds chunk size, split it
            if sentence_length > self.chunk_size:
                # Flush current chunk first
                if current_chunk:
                    chunks.append({
                        "text": " ".join(current_chunk),
                        "source_page": source_page
                    })
                    current_chunk = []
                    current_length = 0
                
                # Split long sentence by words
                words = sentence.split()
                temp_chunk = []
                temp_length = 0
                
                for word in words:
                    if temp_length + len(word) + 1 > self.chunk_size:
                        if temp_chunk:
                            chunks.append({
                                "text": " ".join(temp_chunk),
                                "source_page": source_page
                            })
                        temp_chunk = [word]
                        temp_length = len(word)
                    else:
                        temp_chunk.append(word)
                        temp_length += len(word) + 1
                
                if temp_chunk:
                    current_chunk = temp_chunk
                    current_length = temp_length
                continue
            
            # Check if adding this sentence exceeds chunk size
            if current_length + sentence_length + 1 > self.chunk_size:
                # Save current chunk
                if current_chunk:
                    chunks.append({
                        "text": " ".join(current_chunk),
                        "source_page": source_page
                    })
                
                # Start new chunk with overlap
                if self.overlap > 0 and current_chunk:
                    # Get last few sentences for overlap
                    overlap_text = " ".join(current_chunk)
                    if len(overlap_text) > self.overlap:
                        overlap_text = overlap_text[-self.overlap:]
                        # Find start of sentence in overlap
                        match = re.search(r'[.!?]\s+', overlap_text)
                        if match:
                            overlap_text = overlap_text[match.end():]
                    current_chunk = [overlap_text] if overlap_text.strip() else []
                    current_length = len(overlap_text)
                else:
                    current_chunk = []
                    current_length = 0
            
            current_chunk.append(sentence)
            current_length += sentence_length + 1
        
        # Don't forget the last chunk
        if current_chunk:
            chunks.append({
                "text": " ".join(current_chunk),
                "source_page": source_page
            })
        
        return chunks
    
    # =========================================================================
    # OCR FALLBACK (SLOW PATH - ONLY FOR SCANNED DOCS)
    # =========================================================================
    
    def extract_with_ocr_fallback(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Fall back to OCR for scanned documents.
        Uses unstructured with fast strategy.
        """
        if not self.enable_ocr_fallback:
            print(f"  Skipping OCR fallback (disabled): {file_path.name}")
            return []
        
        try:
            from unstructured.partition.auto import partition
            
            print(f"  Using OCR fallback for scanned document: {file_path.name}")
            self.stats.used_ocr_fallback += 1
            
            elements = partition(
                filename=str(file_path),
                strategy="fast",  # Use fast even for OCR
                include_page_breaks=True
            )
            
            pages = []
            current_page = 1
            current_text = []
            
            for element in elements:
                page_num = getattr(element.metadata, 'page_number', current_page) if hasattr(element, 'metadata') else current_page
                
                if page_num != current_page and current_text:
                    pages.append({
                        "page_number": current_page,
                        "text": "\n".join(current_text),
                        "char_count": sum(len(t) for t in current_text)
                    })
                    current_text = []
                    current_page = page_num
                
                if element.text and element.text.strip():
                    current_text.append(element.text)
            
            # Last page
            if current_text:
                pages.append({
                    "page_number": current_page,
                    "text": "\n".join(current_text),
                    "char_count": sum(len(t) for t in current_text)
                })
            
            return pages
            
        except Exception as e:
            print(f"  OCR fallback failed: {e}")
            return []
    
    # =========================================================================
    # MAIN PROCESSING LOGIC
    # =========================================================================
    
    def process_single_file(self, file_path: Path) -> Optional[List[Dict[str, Any]]]:
        """
        Process a single document and return chunks.
        
        Args:
            file_path: Path to document
            
        Returns:
            List of processed chunks with metadata
        """
        try:
            ext = file_path.suffix.lower()
            file_size_mb = file_path.stat().st_size / 1024 / 1024
            
            print(f"Processing: {file_path.name} ({file_size_mb:.1f} MB)")
            
            # Extract based on file type
            pages = []
            needs_ocr = False
            
            if ext == '.pdf':
                pages, needs_ocr = self.extract_pdf_native(file_path)
                if needs_ocr and self.enable_ocr_fallback:
                    pages = self.extract_with_ocr_fallback(file_path)
            elif ext in ['.docx', '.doc']:
                pages = self.extract_docx(file_path)
            elif ext in ['.pptx', '.ppt']:
                pages = self.extract_pptx(file_path)
            elif ext in ['.txt', '.md']:
                pages = self.extract_text_file(file_path)
            else:
                print(f"  Unsupported format: {ext}")
                return None
            
            if not pages:
                print(f"  No content extracted from {file_path.name}")
                return []
            
            # Generate chunks from extracted pages
            all_chunks = []
            file_hash = self.get_file_hash(file_path)
            
            for page_data in pages:
                page_chunks = self.smart_chunk_text(
                    page_data["text"],
                    source_page=page_data["page_number"]
                )
                
                for i, chunk in enumerate(page_chunks):
                    chunk_data = {
                        "chunk_id": f"{file_path.stem}_p{page_data['page_number']:03d}_c{i:03d}",
                        "source_file": str(file_path.relative_to(self.input_dir)),
                        "file_hash": file_hash,
                        "chunk_index": len(all_chunks),
                        "text": chunk["text"],
                        "metadata": {
                            "filename": file_path.name,
                            "file_directory": str(file_path.parent.relative_to(self.input_dir)) if file_path.parent != self.input_dir else ".",
                            "file_size_bytes": file_path.stat().st_size,
                            "page_number": chunk["source_page"],
                            "processed_timestamp": datetime.now().isoformat(),
                            "extraction_method": "ocr_fallback" if needs_ocr else "native"
                        }
                    }
                    
                    # Only add chunks with meaningful content
                    if chunk_data["text"] and len(chunk_data["text"].strip()) > 10:
                        all_chunks.append(chunk_data)
            
            self.stats.total_pages += len(pages)
            print(f"  Extracted {len(pages)} pages -> {len(all_chunks)} chunks")
            
            return all_chunks
            
        except Exception as e:
            print(f"  Error processing {file_path}: {e}")
            self.failed_files.append({"file": str(file_path), "error": str(e)})
            return None
    
    def process_all_documents(self, parallel: bool = True) -> Dict[str, Any]:
        """
        Process all documents in the input directory.
        
        Args:
            parallel: Whether to use parallel processing
            
        Returns:
            Complete knowledge base dictionary
        """
        start_time = datetime.now()
        documents = self.find_documents()
        self.stats.total_files = len(documents)
        
        print(f"\n{'='*60}")
        print(f"FAST DOCUMENT PROCESSOR")
        print(f"{'='*60}")
        print(f"Input directory: {self.input_dir}")
        print(f"Documents found: {len(documents)}")
        print(f"Workers: {self.workers}")
        print(f"Chunk size: {self.chunk_size} chars")
        print(f"{'='*60}\n")
        
        all_chunks = []
        
        if parallel and len(documents) > 1 and self.workers > 1:
            # Parallel processing
            print(f"Processing {len(documents)} files in parallel with {self.workers} workers...\n")
            
            with ProcessPoolExecutor(max_workers=self.workers) as executor:
                # Submit all tasks
                future_to_path = {
                    executor.submit(self._process_file_wrapper, doc): doc 
                    for doc in documents
                }
                
                # Collect results with progress bar
                for future in tqdm(as_completed(future_to_path), total=len(documents), desc="Processing"):
                    doc_path = future_to_path[future]
                    try:
                        chunks = future.result()
                        if chunks:
                            all_chunks.extend(chunks)
                            self.stats.successful += 1
                            self.stats.total_chunks += len(chunks)
                        elif chunks is not None:
                            self.stats.successful += 1  # Empty but processed
                        else:
                            self.stats.failed += 1
                    except Exception as e:
                        print(f"Error processing {doc_path}: {e}")
                        self.stats.failed += 1
                        self.failed_files.append({"file": str(doc_path), "error": str(e)})
        else:
            # Sequential processing
            print(f"Processing {len(documents)} files sequentially...\n")
            
            for doc_path in tqdm(documents, desc="Processing"):
                chunks = self.process_single_file(doc_path)
                if chunks:
                    all_chunks.extend(chunks)
                    self.stats.successful += 1
                    self.stats.total_chunks += len(chunks)
                elif chunks is not None:
                    self.stats.successful += 1
                else:
                    self.stats.failed += 1
        
        # Calculate timing
        end_time = datetime.now()
        self.stats.processing_time_seconds = (end_time - start_time).total_seconds()
        self.stats.files_per_second = self.stats.total_files / max(self.stats.processing_time_seconds, 0.001)
        
        # Build knowledge base
        knowledge_base = {
            "metadata": {
                "total_documents": self.stats.total_files,
                "successfully_processed": self.stats.successful,
                "failed_documents": self.stats.failed,
                "total_chunks": self.stats.total_chunks,
                "total_pages": self.stats.total_pages,
                "processing_timestamp": datetime.now().isoformat(),
                "processing_time_seconds": round(self.stats.processing_time_seconds, 2),
                "files_per_second": round(self.stats.files_per_second, 2),
                "input_directory": str(self.input_dir),
                "chunk_size": self.chunk_size,
                "overlap": self.overlap,
                "used_ocr_fallback": self.stats.used_ocr_fallback,
                "processor": "fast_processor"
            },
            "chunks": all_chunks,
            "failed_files": self.failed_files,
            "file_index": self._create_file_index(all_chunks)
        }
        
        return knowledge_base
    
    def _process_file_wrapper(self, file_path: Path) -> Optional[List[Dict[str, Any]]]:
        """Wrapper for parallel processing (must be picklable)."""
        return self.process_single_file(file_path)
    
    def _create_file_index(self, chunks: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Create an index of files for quick lookup."""
        file_index = {}
        for chunk in chunks:
            source_file = chunk["source_file"]
            if source_file not in file_index:
                file_index[source_file] = {
                    "chunk_count": 0,
                    "total_characters": 0,
                    "file_directory": chunk["metadata"]["file_directory"],
                    "filename": chunk["metadata"]["filename"],
                    "file_size_bytes": chunk["metadata"]["file_size_bytes"]
                }
            
            file_index[source_file]["chunk_count"] += 1
            file_index[source_file]["total_characters"] += len(chunk["text"])
        
        return file_index
    
    def save_knowledge_base(self, knowledge_base: Dict[str, Any], format_type: str = "both"):
        """Save the processed knowledge base to files."""
        
        if format_type in ["json", "both"]:
            # Full knowledge base with metadata
            json_path = self.output_dir / "knowledge_base.json"
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(knowledge_base, f, indent=2, ensure_ascii=False)
            print(f"Saved knowledge base to: {json_path}")
            
            # Chunks-only JSON file (GPT-friendly - no code interpreter needed)
            chunks_path = self.output_dir / "chunks.json"
            with open(chunks_path, 'w', encoding='utf-8') as f:
                json.dump(knowledge_base["chunks"], f, indent=2, ensure_ascii=False)
            print(f"Saved chunks (GPT-friendly) to: {chunks_path}")
        
        if format_type in ["jsonl", "both"]:
            jsonl_path = self.output_dir / "knowledge_base_chunks.jsonl"
            with open(jsonl_path, 'w', encoding='utf-8') as f:
                for chunk in knowledge_base["chunks"]:
                    f.write(json.dumps(chunk, ensure_ascii=False) + '\n')
            print(f"Saved chunks (JSONL) to: {jsonl_path}")
        
        # Save metadata separately
        metadata_path = self.output_dir / "processing_metadata.json"
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump({
                "metadata": knowledge_base["metadata"],
                "file_index": knowledge_base["file_index"],
                "failed_files": knowledge_base["failed_files"]
            }, f, indent=2, ensure_ascii=False)
        print(f"Saved metadata to: {metadata_path}")
        
        # Create summary CSV
        self._create_summary_csv(knowledge_base)
    
    def _create_summary_csv(self, knowledge_base: Dict[str, Any]):
        """Create a CSV summary of processed files."""
        summary_data = []
        for file_path, info in knowledge_base["file_index"].items():
            summary_data.append({
                "source_file": file_path,
                "filename": info["filename"],
                "directory": info["file_directory"],
                "chunk_count": info["chunk_count"],
                "total_characters": info["total_characters"],
                "file_size_bytes": info["file_size_bytes"]
            })
        
        if summary_data:
            df = pd.DataFrame(summary_data)
            csv_path = self.output_dir / "processing_summary.csv"
            df.to_csv(csv_path, index=False)
            print(f"Saved summary to: {csv_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Fast document processor using native text extraction",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python fast_processor.py --input-dir "src/My Documents" --output-dir "output"
  python fast_processor.py -i "documents" -o "knowledge_base" --workers 4
  python fast_processor.py -i "pdfs" --no-parallel  # Sequential processing
        """
    )
    parser.add_argument("--input-dir", "-i", 
                       required=True,
                       help="Input directory containing documents")
    parser.add_argument("--output-dir", "-o", 
                       default="processed_knowledge_base",
                       help="Output directory for processed files")
    parser.add_argument("--chunk-size", "-c", 
                       type=int, default=1000,
                       help="Maximum characters per chunk (default: 1000)")
    parser.add_argument("--overlap", 
                       type=int, default=200,
                       help="Character overlap between chunks (default: 200)")
    parser.add_argument("--workers", "-w",
                       type=int, default=None,
                       help="Number of parallel workers (default: CPU count - 1)")
    parser.add_argument("--no-parallel",
                       action="store_true",
                       help="Disable parallel processing")
    parser.add_argument("--no-ocr-fallback",
                       action="store_true",
                       help="Disable OCR fallback for scanned documents")
    parser.add_argument("--format", "-f", 
                       choices=["json", "jsonl", "both"], 
                       default="both",
                       help="Output format (default: both)")
    
    args = parser.parse_args()
    
    # Initialize processor
    processor = FastDocumentProcessor(
        input_dir=args.input_dir,
        output_dir=args.output_dir,
        chunk_size=args.chunk_size,
        overlap=args.overlap,
        workers=args.workers,
        enable_ocr_fallback=not args.no_ocr_fallback
    )
    
    # Process documents
    knowledge_base = processor.process_all_documents(parallel=not args.no_parallel)
    
    # Save results
    processor.save_knowledge_base(knowledge_base, args.format)
    
    # Print summary
    meta = knowledge_base['metadata']
    print("\n" + "="*60)
    print("PROCESSING COMPLETE")
    print("="*60)
    print(f"Total documents:     {meta['total_documents']}")
    print(f"Successfully processed: {meta['successfully_processed']}")
    print(f"Failed:              {meta['failed_documents']}")
    print(f"Total pages:         {meta['total_pages']}")
    print(f"Total chunks:        {meta['total_chunks']}")
    print(f"Processing time:     {meta['processing_time_seconds']:.2f} seconds")
    print(f"Speed:               {meta['files_per_second']:.2f} files/second")
    if meta['used_ocr_fallback'] > 0:
        print(f"Used OCR fallback:   {meta['used_ocr_fallback']} files")
    print(f"Output directory:    {args.output_dir}")
    
    if knowledge_base['failed_files']:
        print(f"\nFailed files:")
        for failed in knowledge_base['failed_files']:
            print(f"  - {failed['file']}: {failed['error']}")


if __name__ == "__main__":
    main()
